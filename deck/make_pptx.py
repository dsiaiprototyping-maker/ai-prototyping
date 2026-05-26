#!/usr/bin/env python3
"""
make_pptx.py — DSI May Away Day 2026 · AI Prototyping Workshop
Run: python3 make_pptx.py
Requires: pip install python-pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ── Colours ──────────────────────────────────────────────────────────────────
BG     = RGBColor(0x08, 0x08, 0x0E)
PURPLE = RGBColor(0x7B, 0x2D, 0x8B)
WHITE  = RGBColor(0xDC, 0xDC, 0xE8)
DIM    = RGBColor(0x88, 0x88, 0xA0)
LABEL  = RGBColor(0x6A, 0x6A, 0x88)
CARD   = RGBColor(0x12, 0x08, 0x1A)
DIVCLR = RGBColor(0x2A, 0x10, 0x33)

W = Inches(13.333)
H = Inches(7.5)

# ── Layout constants ──────────────────────────────────────────────────────────
MARGIN    = 0.45   # left/right slide margin (inches)
HDR_RULE  = 0.62   # y of line below chrome row
BODY_TOP  = 1.75   # y where content area starts (below slide title)
TKW_Y     = 6.25   # y of takeaway strip top
CNT_Y     = 7.08   # y of counter
INNER_W   = 13.333 - 2 * MARGIN  # usable width = 12.433"


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def dark_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG


def _run(para, text, size, color, bold=False, italic=False):
    run = para.add_run()
    run.text            = text
    run.font.size       = Pt(size)
    run.font.color.rgb  = color
    run.font.bold       = bold
    run.font.italic     = italic
    return run


def box(slide, text, l, t, w, h,
        size=16, color=WHITE, bold=False, italic=False,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    _run(p, text, size, color, bold, italic)
    return txb, tf


def bullets(slide, items, l, t, w, h,
            size=17, color=WHITE, prefix='·  ', spacing=3):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing)
        _run(p, prefix + item, size, color)
    return txb


def rect(slide, x, y, w, h, fill_color, line_color=None, line_w=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_w:
            shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def vdiv(slide, x, y_top=BODY_TOP, y_bot=TKW_Y):
    rect(slide, x, y_top, 0.012, y_bot - y_top, DIVCLR)


def chrome(slide, center):
    box(slide, 'DSI · AWAY DAY', MARGIN, 0.16, 3.5, 0.38, size=9,  color=LABEL)
    box(slide, center,            3.5,   0.16, 6.5, 0.38, size=9,  color=LABEL, align=PP_ALIGN.CENTER)
    box(slide, 'MAY 2026',        9.5,   0.16, 3.3, 0.38, size=9,  color=LABEL, align=PP_ALIGN.RIGHT)
    rect(slide, 0, HDR_RULE, 13.333, 0.01, DIVCLR)


def slabel(slide, text, y=0.70):
    box(slide, text, MARGIN, y, INNER_W, 0.34, size=12, color=PURPLE)


def stitle(slide, text, y=0.96):
    box(slide, text, MARGIN, y, INNER_W, 0.72, size=30, color=WHITE)
    rect(slide, 0, BODY_TOP, 13.333, 0.01, DIVCLR)


def takeaway(slide, text):
    rect(slide, 0, TKW_Y - 0.06, 13.333, 0.82, RGBColor(0x10, 0x05, 0x16))
    box(slide, text, MARGIN, TKW_Y, INNER_W, 0.72,
        size=17, color=DIM, italic=True)


def counter(slide, n, total=5):
    box(slide, f'{n:02d} / {total:02d}', 10.5, CNT_Y, 2.5, 0.4,
        size=13, color=PURPLE, align=PP_ALIGN.RIGHT)


# ── SLIDE 1: COVER ────────────────────────────────────────────────────────────
def make_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    chrome(slide, 'SLIDE 01 · WELCOME')

    # Eyebrow
    box(slide, 'DSI · MAY AWAY DAY · 2026',
        0.6, 1.85, 10, 0.48, size=13, color=PURPLE)

    # Main headline — massive, fills the slide horizontally
    box(slide, 'AI Prototyping Workshop',
        0.6, 2.35, 12.2, 1.55, size=56, color=WHITE)

    # Subtitle
    box(slide, 'From prompt to shipped — live.',
        0.6, 3.9, 10, 0.7, size=23, color=DIM, italic=True)

    # Purple rule
    rect(slide, 0.6, 4.82, 0.65, 0.03, PURPLE)

    # NatWest brand
    box(slide, 'NatWest',
        0.6, 5.06, 5, 0.72, size=30, color=PURPLE, bold=True)
    box(slide, 'DATA SCIENCE & INNOVATION',
        0.6, 5.72, 8, 0.4, size=11, color=LABEL)

    counter(slide, 1)


# ── SLIDE 2: LLM vs AGENT ─────────────────────────────────────────────────────
def make_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    chrome(slide, 'SLIDE 02 · LLM CHAT VS AI AGENT HARNESS')
    slabel(slide, 'FOUNDATIONS')
    stitle(slide, 'LLM Chat vs AI Agent Harness')

    # Each column is half of usable width; split at x=6.667
    COL_W = INNER_W / 2 - 0.25  # ~5.97"
    LX = MARGIN          # left col x
    RX = 13.333 / 2 + 0.2  # right col x

    # Column sub-headers
    box(slide, 'LLM / CHAT INTERFACE', LX, 1.82, COL_W, 0.42, size=12, color=LABEL)
    box(slide, 'AI AGENT HARNESS',     RX, 1.82, COL_W, 0.42, size=12, color=PURPLE)

    # Horizontal rule under headers
    rect(slide, 0, 2.3, 13.333, 0.008, DIVCLR)

    left = [
        'You ask a question, it gives an answer',
        'Mostly conversation',
        'Best for brainstorming, writing, summarising',
        'Human does most of the execution',
        'Output is usually text',
    ]
    right = [
        'You give a goal, it can take actions to achieve it',
        'Conversation plus tools, memory, permissions, and workflows',
        'Best for building, checking, testing, fixing, and iterating',
        'Agent can do parts of the execution under supervision',
        'Output can be code, edits, screenshots, test results, PRs, browser actions',
    ]

    # Bullets fill from 2.38 to TKW_Y
    BH = TKW_Y - 2.42 - 0.05   # ≈3.78"
    bullets(slide, left,  LX, 2.42, COL_W, BH, size=17, color=DIM,   spacing=6)
    bullets(slide, right, RX, 2.42, COL_W, BH, size=17, color=WHITE, spacing=6)

    # Vertical centre divider
    vdiv(slide, 13.333 / 2 - 0.006, 1.75)

    takeaway(slide, 'A chat interface gives you answers. An agent harness helps turn answers into completed work.')
    counter(slide, 2)


# ── SLIDE 3: THREE WAYS ───────────────────────────────────────────────────────
def make_slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    chrome(slide, 'SLIDE 03 · CLAUDE DESKTOP APP — THREE WAYS TO WORK')
    slabel(slide, 'TOOLS')
    stitle(slide, 'Claude Desktop App — Three Ways to Work')

    cols = [
        ('CHAT', 'Stateless', [
            'One-off conversation',
            'No memory between sessions',
            'Best for questions and drafts',
            'Output is text',
            'You drive everything',
            'Great for quick thinking, summarising, explaining',
        ]),
        ('PROJECTS (COWORK)', 'Stateful', [
            'Persistent workspace',
            'Shared context across sessions',
            'Best for ongoing team work',
            'Output is text plus files',
            'Claude remembers your project',
            'Great for collaborations spanning days or weeks',
        ]),
        ('CLAUDE CODE', 'Agentic', [
            'Agent inside your codebase',
            'Reads, edits, and runs files',
            'Best for building and iterating',
            'Output is code, tests, PRs, terminal actions',
            'Claude acts with your supervision',
            'Great for software tasks where Claude sees and changes artefacts',
        ]),
    ]

    N = len(cols)
    gap = 0.18
    cw = (INNER_W - gap * (N - 1)) / N   # ≈4.03"
    BH = TKW_Y - 2.78 - 0.05             # ≈3.42" for bullets

    for i, (tag, title, items) in enumerate(cols):
        x = MARGIN + i * (cw + gap)
        box(slide, tag,   x, 1.82, cw, 0.38, size=12, color=PURPLE)
        box(slide, title, x, 2.22, cw, 0.54, size=22, color=WHITE)
        rect(slide, x, 2.78, cw, 0.008, DIVCLR)
        bullets(slide, items, x, 2.80, cw, BH, size=16, color=DIM, spacing=5)
        if i < N - 1:
            vdiv(slide, x + cw + gap / 2, 1.75)

    takeaway(slide, 'Chat answers your question. Projects remembers your work. Claude Code does the work with you.')
    counter(slide, 3)


# ── SLIDE 4: SETTING UP ───────────────────────────────────────────────────────
def make_slide4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    chrome(slide, 'SLIDE 04 · SETTING CLAUDE UP FOR SUCCESS')
    slabel(slide, 'CONFIGURATION')
    stitle(slide, 'Setting Claude Up for Success')

    levers = [
        ('01 · GIT', 'Version Control',
         'Gives Claude a safety net — diff, revert, and branch. '
         'Commit history is context Claude can read and reason about. '
         'Branches let Claude work in isolation before you review.'),
        ('02 · MEMORY', 'Persistent Facts',
         'Persistent facts across sessions. Stored as markdown files '
         'Claude reads at start of conversation. Types: user profile, '
         'project state, feedback corrections, external references.'),
        ('03 · CLAUDE.MD', 'Rules File',
         'Project-level instructions checked into the repo. Define stack '
         'constraints, naming conventions, scope limits. Claude treats this '
         'as the source of truth for how to work in the project.'),
        ('04 · SKILLS', 'Slash Commands',
         'Reusable capabilities invoked with a slash command. '
         'Examples: /review, /plan, /security-review, /init, /graphify. '
         'Custom skills can encode team-specific workflows.'),
        ('05 · PERMISSIONS', 'Allowlist',
         'Explicit allowlist of what Claude can run automatically. '
         'Reduces friction for low-risk commands (read, test, lint). '
         'High-risk actions (push, deploy, delete) still prompt for approval.'),
    ]

    # Layout: row 1 → 3 equal cards; row 2 → 2 wider cards (full width)
    VGAP   = 0.16
    HGAP3  = 0.14   # gap between 3 cards
    HGAP2  = 0.18   # gap between 2 cards

    R1_Y   = BODY_TOP + 0.08
    R1_H   = 2.22
    R2_Y   = R1_Y + R1_H + VGAP
    R2_H   = TKW_Y - R2_Y - 0.12

    CW3    = (INNER_W - HGAP3 * 2) / 3    # ≈4.07"
    CW2    = (INNER_W - HGAP2)     / 2    # ≈6.13"

    # Row 1 — three cards
    for i in range(3):
        x = MARGIN + i * (CW3 + HGAP3)
        num, title, body = levers[i]
        rect(slide, x, R1_Y, CW3, R1_H, CARD,
             line_color=RGBColor(0x4A, 0x1A, 0x55), line_w=0.6)
        box(slide, num,   x+0.16, R1_Y+0.14, CW3-0.32, 0.36, size=11, color=PURPLE)
        box(slide, title, x+0.16, R1_Y+0.46, CW3-0.32, 0.52, size=19, color=WHITE)
        box(slide, body,  x+0.16, R1_Y+0.96, CW3-0.32, R1_H-1.08,
            size=15, color=DIM, wrap=True)

    # Row 2 — two wider cards
    for i in range(2):
        x = MARGIN + i * (CW2 + HGAP2)
        num, title, body = levers[3 + i]
        rect(slide, x, R2_Y, CW2, R2_H, CARD,
             line_color=RGBColor(0x4A, 0x1A, 0x55), line_w=0.6)
        box(slide, num,   x+0.18, R2_Y+0.14, CW2-0.36, 0.36, size=11, color=PURPLE)
        box(slide, title, x+0.18, R2_Y+0.46, CW2-0.36, 0.52, size=19, color=WHITE)
        box(slide, body,  x+0.18, R2_Y+0.96, CW2-0.36, R2_H-1.08,
            size=15, color=DIM, wrap=True)

    takeaway(slide, 'The more context and guardrails you give Claude upfront, the less supervision it needs mid-task.')
    counter(slide, 4)


# ── SLIDE 5: BUILDING — WORKFLOW INFOGRAPHIC ──────────────────────────────────
def make_slide5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    chrome(slide, 'SLIDE 05 · BUILDING WITH CLAUDE — FROM PROMPT TO SHIPPED')
    slabel(slide, 'PROCESS')
    stitle(slide, 'Building with Claude — From Prompt to Shipped')

    phases = [
        ('01', 'Good Prompting', [
            'Point at what you see',
            'Name the constraint',
            'One visible goal per prompt',
        ]),
        ('02', 'Plan Phase', [
            'Ask Claude to plan first',
            'Read and reject wrong assumptions',
            'Opus for complex decisions',
        ]),
        ('03', 'Interview Phase', [
            'Claude asks — answer them',
            'Hidden conflicts surface here',
            'Shared understanding saves rework',
        ]),
        ('04', 'Screenshot Feedback', [
            'Capture what you see',
            'Paste and describe the issue',
            'Claude reads and corrects visually',
        ]),
        ('05', 'Iteration Loop', [
            'prompt → edit → refresh → check',
            'Small prompts = short diffs',
            'Commit after each working step',
        ]),
    ]

    N        = len(phases)
    SLOT_W   = INNER_W / N          # 12.433/5 = 2.4866"
    NODE_D   = 0.80                  # oval diameter
    NODE_Y   = 1.92                  # top of node row
    NODE_CY  = NODE_Y + NODE_D / 2   # 2.32"
    TITLE_Y  = NODE_Y + NODE_D + 0.15  # 2.87"
    TITLE_H  = 0.62
    BULLETS_Y = TITLE_Y + TITLE_H    # 3.49"
    LOOP_Y   = 5.55
    BH       = LOOP_Y - BULLETS_Y - 0.10  # ≈2.06"
    SLOT_PAD = 0.10
    ARROW_Y  = NODE_CY - 0.10       # vertically centred on node
    ARROW_H  = 0.20
    ARROW_W  = SLOT_W - NODE_D       # gap between node edges ≈1.687"

    for i, (num, title, items) in enumerate(phases):
        cx = MARGIN + (i + 0.5) * SLOT_W
        nx = cx - NODE_D / 2

        # Oval node
        node = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(nx), Inches(NODE_Y), Inches(NODE_D), Inches(NODE_D)
        )
        is_last = (i == N - 1)
        node.fill.solid()
        node.fill.fore_color.rgb = (
            RGBColor(0x2A, 0x0A, 0x38) if is_last else RGBColor(0x14, 0x06, 0x1C)
        )
        node.line.color.rgb = PURPLE
        node.line.width = Pt(2.0 if is_last else 1.2)

        # Phase number centred inside oval
        box(slide, num,
            nx, NODE_Y + NODE_D * 0.22, NODE_D, NODE_D * 0.56,
            size=15, color=PURPLE, align=PP_ALIGN.CENTER)

        # Phase title below node
        tx = MARGIN + i * SLOT_W + SLOT_PAD
        tw = SLOT_W - 2 * SLOT_PAD
        box(slide, title, tx, TITLE_Y, tw, TITLE_H,
            size=14, color=WHITE, italic=True, align=PP_ALIGN.CENTER, wrap=True)

        # Bullet points
        bullets(slide, items, tx, BULLETS_Y, tw, BH,
                size=12, color=DIM, prefix='· ', spacing=5)

        # Right-arrow connector to next phase
        if i < N - 1:
            ax = nx + NODE_D
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                Inches(ax), Inches(ARROW_Y), Inches(ARROW_W), Inches(ARROW_H)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = PURPLE
            arrow.line.fill.background()

    # Loop bar
    rect(slide, MARGIN, LOOP_Y, INNER_W, 0.012, DIVCLR)
    box(slide, '↺   LOOP  ·  ITERATE  ·  SHIP',
        MARGIN, LOOP_Y + 0.02, INNER_W, 0.28,
        size=10, color=LABEL, align=PP_ALIGN.CENTER)
    rect(slide, MARGIN, LOOP_Y + 0.32, INNER_W, 0.012, DIVCLR)

    takeaway(slide, 'Build in short loops: prompt, look, screenshot, correct — the browser is your shared source of truth.')
    counter(slide, 5)


# ── MAIN ──────────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    prs = new_prs()
    make_slide1(prs)
    make_slide2(prs)
    make_slide3(prs)
    make_slide4(prs)
    make_slide5(prs)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'presentation.pptx')
    prs.save(out)
    print(f'Saved → {out}')
